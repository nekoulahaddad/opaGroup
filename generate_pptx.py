from pathlib import Path
import tempfile
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

base = Path(__file__).resolve().parent
assets = base / "assets"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

brand = RGBColor(3, 46, 130)
muted = RGBColor(71, 85, 105)

temp_files = []

def add_title(slide, title, subtitle=None):
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.8), Inches(1.2))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = brand
    if subtitle:
        p2 = tx.text_frame.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = muted


def add_image(slide, name, x, y, w, h):
    path = assets / name
    if not path.exists():
        return
    use_path = path
    if path.suffix.lower() == ".webp":
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
        tmp.close()
        with PILImage.open(path) as im:
            im = im.convert("RGB")
            im.save(tmp.name, format="PNG")
        temp_files.append(tmp.name)
        use_path = Path(tmp.name)
    slide.shapes.add_picture(str(use_path), x, y, w, h)


def add_link_paragraph(frame, label, url, text):
    p = frame.add_paragraph() if frame.paragraphs and frame.paragraphs[0].text else frame.paragraphs[0]
    p.text = ""
    run = p.add_run()
    run.text = f"{label}: {text}"
    run.font.size = Pt(16)
    run.hyperlink.address = url

# Cover
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_image(slide, "hero.webp", Inches(0), Inches(0), prs.slide_width, prs.slide_height)
box = slide.shapes.add_textbox(Inches(0.9), Inches(0.7), Inches(11.5), Inches(2.2))
frame = box.text_frame
p = frame.paragraphs[0]
p.text = "OPA GROUP CO LTD"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
sp = frame.add_paragraph()
sp.text = "International Distribution, Logistics & Trade"
sp.font.size = Pt(20)
sp.font.color.rgb = RGBColor(255, 255, 255)
sp2 = frame.add_paragraph()
run = sp2.add_run()
run.text = "opagroupafrica.site"
run.font.size = Pt(16)
run.font.color.rgb = RGBColor(255, 255, 255)
run.hyperlink.address = "https://opagroupafrica.site/"

# Company Overview
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Company Overview", "Who we are and what we deliver")
add_image(slide, "page6_img2.jpg", Inches(8.2), Inches(1.6), Inches(4.5), Inches(3.2))
add_image(slide, "page1_img4.jpg", Inches(0.8), Inches(4.6), Inches(3.8), Inches(2.6))
add_image(slide, "page1_img5.jpg", Inches(4.9), Inches(4.6), Inches(3.2), Inches(2.6))
text = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(7.2), Inches(2.8))
frame = text.text_frame
frame.word_wrap = True
p = frame.paragraphs[0]
p.text = (
    "OPA Group Co. Ltd is a dynamic company specializing in distribution, logistics, and supply chain management across Africa, Asia, and the Middle East. "
    "We connect manufacturers and suppliers directly with regional distributors, ensuring timely access to quality products and streamlined trade operations."
)
p.font.size = Pt(16)

# Leadership
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Leadership & Strategy", "Executive direction and growth priorities")
add_image(slide, "page5_img2.jpg", Inches(8.2), Inches(1.6), Inches(4.5), Inches(5.5))
text = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(7.2), Inches(5.6))
frame = text.text_frame
frame.word_wrap = True
p = frame.paragraphs[0]
p.text = (
    "OPA Group has grown from a regional supplier into a multinational logistics and distribution powerhouse operating across 12+ countries. "
    "We invest in technology, partnerships, and skilled people to deliver transparent, efficient, and world-class services."
)
p.font.size = Pt(16)

# Core Services
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Core Services")
add_image(slide, "page10_img1.jpg", Inches(0.8), Inches(4.2), Inches(4.0), Inches(2.6))
add_image(slide, "page11_img1.jpg", Inches(5.1), Inches(4.2), Inches(4.0), Inches(2.6))
add_image(slide, "page11_img2.jpg", Inches(9.4), Inches(4.2), Inches(3.6), Inches(2.6))
text = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(12.0), Inches(2.2))
frame = text.text_frame
frame.word_wrap = True
for i, t in enumerate([
    "Distribution & Wholesale Supply",
    "Logistics & Transportation",
    "Import & Export Services",
    "Supply Chain Management",
    "Procurement & Sourcing",
    "Customs Clearance & Freight Forwarding",
    "Specialized Trading Solutions",
    "Shipping & Clearing",
    "Market Support",
]):
    p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
    p.text = f"• {t}"
    p.font.size = Pt(16)

# Market Support
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Market Support & Extended Solutions")
add_image(slide, "page12_img3.jpg", Inches(0.8), Inches(4.2), Inches(4.0), Inches(2.6))
add_image(slide, "page13_img1.jpg", Inches(5.1), Inches(4.2), Inches(4.0), Inches(2.6))
add_image(slide, "page13_img2.jpg", Inches(9.4), Inches(4.2), Inches(3.6), Inches(2.6))
text = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(12.0), Inches(2.4))
frame = text.text_frame
p = frame.paragraphs[0]
p.text = (
    "We build sales and marketing teams, develop distribution channels, and accelerate brand performance. "
    "We also provide electronics supply, cement and construction materials, vehicle sales & rental, apartments sales & rental, and hotel & flight booking services."
)
p.font.size = Pt(16)

# Products
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Products & Market Categories")
add_image(slide, "page9_img1.jpg", Inches(8.2), Inches(1.6), Inches(4.5), Inches(5.6))
text = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(7.0), Inches(5.6))
frame = text.text_frame
for i, t in enumerate([
    "Baked goods and breakfast cereals",
    "Dairy products including milk, cheese, butter, yogurt",
    "Beverages: soft drinks, bottled water, coffee, and alcoholic beverages",
    "Fruits and vegetables for local and export markets",
    "Processed and fresh foods, packaged, canned, and frozen",
    "Healthcare products and essential medicines",
    "Personal care and cosmetics",
    "Household and cleaning items",
    "Electronics and technology products",
    "Cement and construction materials",
]):
    p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
    p.text = f"• {t}"
    p.font.size = Pt(16)

# Process
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Operational Footprint")
add_image(slide, "page15_img4.png", Inches(8.2), Inches(1.6), Inches(4.5), Inches(5.6))
text = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(7.0), Inches(5.6))
frame = text.text_frame
for i, t in enumerate([
    "Understanding client needs and market dynamics",
    "Strategic sourcing and procurement from trusted manufacturers",
    "Warehousing and inventory management across East and Southern Africa",
    "Reliable logistics and distribution with a dedicated fleet",
    "Technology-driven operations for tracking and reporting",
    "Quality control and compliance for every shipment",
    "Continuous improvement for long-term partnerships",
]):
    p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
    p.text = f"• {t}"
    p.font.size = Pt(16)

# Presence
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Geographical Presence")
add_image(slide, "page17_img2.jpg", Inches(8.2), Inches(1.6), Inches(4.5), Inches(5.6))
text = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(7.0), Inches(5.6))
frame = text.text_frame
p = frame.paragraphs[0]
p.text = "A growing network across Africa, Asia, and the Middle East with warehouses, distribution hubs, and transport fleets near key trade routes and ports."
p.font.size = Pt(16)

# Contact
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Contact")
text = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(6.0), Inches(5.6))
frame = text.text_frame
add_link_paragraph(frame, "Website", "https://opagroupafrica.site/", "opagroupafrica.site")
add_link_paragraph(frame, "Email", "mailto:info@opagroupafrica.com", "info@opagroupafrica.com")
add_link_paragraph(frame, "Email", "mailto:opagroup66@gmail.com", "opagroup66@gmail.com")
add_link_paragraph(frame, "Email", "mailto:ceo@opagroupafrica.com", "ceo@opagroupafrica.com")
add_link_paragraph(frame, "Phone", "tel:+27842600556", "+27 84 260 0556")
add_link_paragraph(frame, "WhatsApp", "https://wa.me/254758896370", "+254 758 896 370")
add_link_paragraph(frame, "Facebook", "https://www.facebook.com/opagroupcompamy", "OPA GROUP CO. LTD")
add_link_paragraph(frame, "TikTok", "https://www.tiktok.com/@opa.group.co.ltd", "@opa.group.co.ltd")
add_link_paragraph(frame, "Instagram", "https://www.instagram.com/opa_group_co_ltd/", "@opa_group_co_ltd")
add_link_paragraph(frame, "LinkedIn", "https://www.linkedin.com/in/opa-group-co-ltd-251bb3396/", "OPA GROUP CO. LTD")
add_image(slide, "contact2.webp", Inches(7.2), Inches(1.6), Inches(5.8), Inches(5.6))

out = base / "OPA_GROUP_CO_LTD_Presentation.pptx"
prs.save(out)
print(f"Generated: {out}")

for tmp in temp_files:
    try:
        Path(tmp).unlink()
    except OSError:
        pass
